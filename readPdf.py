# import Presentation class 
# from pptx library 
from pptx import Presentation 
  
# creating an object 
ppt = Presentation("Samples/input1.pptx") 
  
# open file in write mode 
File_to_write_data = open("File_To_Extract_ppt.txt", "w") 
  
# write text from powerpoint 
# file into .txt file 
for slide in ppt.slides:  
    for shape in slide.shapes:  
        if not shape.has_text_frame:  
            continue 
        for paragraph in shape.text_frame.paragraphs:  
            for run in paragraph.runs:  
                # File_to_write_data.writelines(run.text+"\n") 
                break
  
# close the file                
File_to_write_data.close() 
  
print("Done")

# text recognition
import numpy as np
# import cv2
